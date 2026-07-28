"""
WHAT THIS FILE IS
-----------------
Generates the Brink final-presentation PowerPoint (brink-deck.pptx). Run it to
(re)build the deck:

    python docs/presentation/build_deck.py

Every slide carries its speaker script in the PowerPoint "notes" pane, so each
presenter has their lines. Styling matches the product's palette (dark canvas,
lavender + hot-pink accents). Several slides show real app screenshots, and the
architecture slide is drawn as a labelled diagram. Edit the build() function and
re-run to change the deck.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# Brand palette (from backend/app/static/brink.css :root)
INK = RGBColor(0x0B, 0x0B, 0x12)      # page background
DEEP = RGBColor(0x10, 0x10, 0x1A)     # slightly lifted panel (diagram container)
PANEL = RGBColor(0x15, 0x15, 0x1F)    # card background
LINE = RGBColor(0x2E, 0x2E, 0x44)     # borders
MUTE = RGBColor(0x8A, 0x8A, 0xA3)     # secondary text
TEXT = RGBColor(0xED, 0xED, 0xF5)     # primary text
ACCENT = RGBColor(0x9D, 0x8D, 0xF1)   # lavender
HOT = RGBColor(0xF4, 0x72, 0xB6)      # hot pink
TEAL = RGBColor(0x5E, 0xEA, 0xD4)     # supporting

FONT = "Inter"  # PowerPoint gracefully substitutes if Inter isn't installed
SHOTS = Path(__file__).parent.parent / "screenshots"


# ---------- low-level helpers ----------

def _set_bg(slide, color=INK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _style(run, size, color=TEXT, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _rule(slide, left, top, width, color=ACCENT, height=Pt(3)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _kicker_title(slide, kicker, title, title_color=ACCENT, title_w=11.9):
    _, ktf = _box(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4))
    kr = ktf.paragraphs[0].add_run(); kr.text = kicker; _style(kr, 13, MUTE, bold=True)
    _, ttf = _box(slide, Inches(0.7), Inches(0.95), Inches(title_w), Inches(1.4))
    for i, line in enumerate(title.split("\n")):
        p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        r = p.add_run(); r.text = line; _style(r, 32, title_color, bold=True)
    _rule(slide, Inches(0.72), Inches(2.2), Inches(1.5))


def _bullets(slide, bullets, left, top, width, height, size=20):
    _, btf = _box(slide, left, top, width, height)
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(14)
        lead = b.startswith("**")
        txt = b.replace("**", "")
        run = p.add_run()
        run.text = ("•  " if not lead else "") + txt
        _style(run, size, TEXT if not lead else HOT, bold=lead)


def _fit_image(slide, path, bx, by, bw, bh):
    from PIL import Image
    try:
        with Image.open(path) as im:
            iw, ih = im.size
        ratio = min(bw / iw, bh / ih)
        w, h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
    except Exception:
        w, h = bw, None
    left = Emu(int(bx + (bw - (w if isinstance(w, Emu) else bw)) / 2))
    top = Emu(int(by + (bh - (h if h else bh)) / 2)) if h else by
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)


# ---------- slide builders ----------

def add_content(prs, kicker, title, bullets, notes, title_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _kicker_title(slide, kicker, title, title_color)
    _bullets(slide, bullets, Inches(0.7), Inches(2.6), Inches(11.9), Inches(4.4))
    _notes(slide, notes)
    return slide


def add_split(prs, kicker, title, bullets, image, notes, title_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _kicker_title(slide, kicker, title, title_color, title_w=7.0)
    _bullets(slide, bullets, Inches(0.7), Inches(2.6), Inches(6.3), Inches(4.4), size=19)
    # framed image panel on the right
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.35), Inches(1.0), Inches(5.3), Inches(5.9))
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE; panel.line.width = Pt(1); panel.shadow.inherit = False
    if image and (SHOTS / image).exists():
        _fit_image(slide, SHOTS / image, Inches(7.5), Inches(1.15), Inches(5.0), Inches(5.6))
    _notes(slide, notes)
    return slide


def add_bignum(prs, kicker, title, stats, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _, ktf = _box(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4))
    kr = ktf.paragraphs[0].add_run(); kr.text = kicker; _style(kr, 13, MUTE, bold=True)
    _, ttf = _box(slide, Inches(0.7), Inches(0.95), Inches(11.9), Inches(1.0))
    tr = ttf.paragraphs[0].add_run(); tr.text = title; _style(tr, 30, TEXT, bold=True)
    for i, (num, label) in enumerate(stats):
        left = Inches(0.7) + i * Inches(4.0)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.7), Inches(3.7), Inches(3.0))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
        _, ntf = _box(slide, left, Inches(3.1), Inches(3.7), Inches(1.5))
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        nr = ntf.paragraphs[0].add_run(); nr.text = num
        _style(nr, 58, ACCENT if i % 2 == 0 else HOT, bold=True)
        _, ltf = _box(slide, left + Inches(0.2), Inches(4.55), Inches(3.3), Inches(1.0))
        ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
        lr = ltf.paragraphs[0].add_run(); lr.text = label; _style(lr, 17, MUTE)
    _, stf = _box(slide, Inches(0.7), Inches(6.1), Inches(11.9), Inches(1.0))
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sr = stf.paragraphs[0].add_run(); sr.text = subtitle; _style(sr, 20, TEXT, italic=True)
    _notes(slide, notes)
    return slide


def add_columns(prs, kicker, title, cols, notes):
    """Three labelled cards side by side (heading + a few lines each)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _kicker_title(slide, kicker, title)
    for i, (head, lines, color) in enumerate(cols):
        left = Inches(0.7) + i * Inches(4.05)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.7), Inches(3.75), Inches(3.5))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = color; card.line.width = Pt(1.5); card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.25)
        hp = tf.paragraphs[0]; hr = hp.add_run(); hr.text = head; _style(hr, 20, color, bold=True)
        hp.space_after = Pt(10)
        for ln in lines:
            p = tf.add_paragraph(); p.space_after = Pt(8)
            r = p.add_run(); r.text = ln; _style(r, 15, TEXT)
    _notes(slide, notes)
    return slide


def add_image(prs, img_path, caption, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _, ctf = _box(slide, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.6))
    cr = ctf.paragraphs[0].add_run(); cr.text = caption; _style(cr, 18, ACCENT, bold=True)
    _fit_image(slide, img_path, Inches(0.7), Inches(1.1), Inches(11.9), Inches(5.9))
    _notes(slide, notes)
    return slide


def add_title(prs, title, subtitle, footer, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _rule(slide, Inches(0.9), Inches(2.55), Inches(2.4), HOT, Pt(5))
    _, ttf = _box(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6))
    tr = ttf.paragraphs[0].add_run(); tr.text = title; _style(tr, 50, TEXT, bold=True)
    _, stf = _box(slide, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.0))
    sr = stf.paragraphs[0].add_run(); sr.text = subtitle; _style(sr, 24, ACCENT, bold=True)
    _, ftf = _box(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.4))
    for i, line in enumerate(footer):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        r = p.add_run(); r.text = line; _style(r, 16, MUTE)
    _notes(slide, notes)
    return slide


# ---------- architecture diagram ----------

def _diag_box(slide, x, y, w, h, title, sub=None, fill=PANEL, edge=LINE, tcolor=TEXT, tsize=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge; shp.line.width = Pt(1.5); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; _style(r, tsize, tcolor, bold=True)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; _style(r2, tsize - 3, MUTE)
    return shp


def _arrow(slide, x1, y1, x2, y2, color=MUTE, width=1.75, both=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle'}))
    if both:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle'}))
    return c


def _mini_label(slide, x, y, w, text, color=MUTE, size=10):
    _, tf = _box(slide, Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = text; _style(r, size, color)


def add_architecture(prs, kicker, title, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _set_bg(slide)
    _, ktf = _box(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.4))
    kr = ktf.paragraphs[0].add_run(); kr.text = kicker; _style(kr, 13, MUTE, bold=True)
    _, ttf = _box(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7))
    tr = ttf.paragraphs[0].add_run(); tr.text = title; _style(tr, 30, ACCENT, bold=True)

    # top row: Spotify (external) + Browser
    _diag_box(slide, 0.7, 1.9, 2.7, 1.0, "Spotify", "logins · listening · catalog", edge=HOT, tsize=13)
    _diag_box(slide, 5.15, 1.9, 3.0, 1.0, "Your browser", "web pages + light JS", edge=ACCENT, tsize=13)

    # centre: the one FastAPI app container with 4 inner components
    _diag_box(slide, 3.15, 3.35, 7.0, 2.0, "FastAPI app  ·  Render", fill=DEEP, edge=ACCENT, tsize=15)
    # nudge the container title to the top so inner boxes sit below it
    cont_title = slide.shapes[-1].text_frame
    cont_title.vertical_anchor = MSO_ANCHOR.TOP
    inner = [("Web pages", "Jinja"), ("JSON API", None), ("Auth", "JWT · AES-256"), ("Inference", "on read")]
    for i, (t, s) in enumerate(inner):
        _diag_box(slide, 3.35 + i * 1.68, 3.95, 1.5, 1.15, t, s, fill=PANEL, edge=LINE, tsize=12)

    # bottom row: Supabase + Analytics
    _diag_box(slide, 2.2, 5.75, 4.6, 1.4, "Supabase", fill=DEEP, edge=TEAL, tsize=14)
    slide.shapes[-1].text_frame.vertical_anchor = MSO_ANCHOR.TOP
    sub = [("Postgres", "public + bronze/silver/gold"), ("Auth", None), ("Storage", None)]
    widths = [2.1, 1.0, 1.0]
    xoff = 2.35
    for (t, s), wd in zip(sub, widths):
        _diag_box(slide, xoff, 6.25, wd - 0.1, 0.8, t, s, fill=PANEL, edge=LINE, tsize=11)
        xoff += wd
    _diag_box(slide, 7.35, 5.75, 3.3, 1.4, "Analytics pipeline",
              "GitHub Actions · nightly\ntrains K-means → writes gold", edge=HOT, tsize=13)

    # arrows
    _arrow(slide, 6.65, 2.9, 6.65, 3.35, color=ACCENT)            # browser -> app
    _mini_label(slide, 6.75, 2.95, 2.4, "HTTPS · same-origin", ACCENT)
    _arrow(slide, 2.6, 2.9, 4.0, 3.35, color=HOT, both=True)      # spotify <-> app
    _mini_label(slide, 0.7, 3.0, 2.3, "OAuth · now playing", HOT)
    _arrow(slide, 5.0, 5.35, 4.6, 5.75, color=TEAL)              # app -> supabase
    _mini_label(slide, 4.9, 5.42, 2.2, "SQLModel ORM", TEAL)
    _arrow(slide, 7.35, 6.45, 6.85, 6.55, color=HOT)            # analytics -> postgres
    _mini_label(slide, 6.7, 6.85, 2.2, "writes gold", HOT)

    # caption
    _, cap = _box(slide, Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.4))
    cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    cr = cap.paragraphs[0].add_run()
    cr.text = "One app is everything the user touches. A nightly job trains the model on the side; the app just reads the results."
    _style(cr, 13, MUTE, italic=True)
    _notes(slide, notes)
    return slide


# ---------- the deck ----------

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 — Title
    add_title(
        prs, "Brink — Music, but social. 🎶",
        "A social app built on what you actually listen to.",
        ["brink-xg7p.onrender.com",
         "Andrea Vreugdenhil (backend) · Jonah Walker (analytics) · Sebastian Arguedas Soley (frontend)",
         "McGill Desautels — MMA"],
        "ANDREA (0:20): Hi everyone — this is Brink. It's a social app built around the music you actually "
        "listen to on Spotify. I worked on the backend, Jonah did the analytics, and Sebastian did the "
        "frontend. It's deployed and running, so we'll show you the real thing at the end.",
    )

    # 2 — The idea (humanised)
    add_content(
        prs, "ANDREA · 1:00", "The idea",
        ["On most apps, what you post is the highlight reel — a version of yourself.",
         "What you listen to is more honest: it's just what you played.",
         "**So we built the feed on that — and used it to find who listens like you.",
         "Two questions run through the whole app: what do you listen to, and who matches you?"],
        "ANDREA (1:00): The starting point was simple. On most social apps, what you post is a choice — you "
        "show the version of yourself you want people to see. Your listening history is different: it's "
        "just what you actually played. We thought that was a more honest thing to build a feed around. So "
        "on Brink you share the songs you're really listening to, people react and comment, and then we "
        "use that listening data to figure out whose taste actually lines up with yours.",
    )

    # 3 — What Brink is (split + feed screenshot)
    add_split(
        prs, "ANDREA · 1:00", "What Brink is",
        ["**Feed — songs that play right in the card, with reactions & comments",
         "**Taste communities — a group our model puts you in",
         "**Compatibility — how close two people's taste is",
         "**Artist portal — behind-the-scenes posts to fans",
         "It's all one live app — everything we planned is built."],
        "Feed-brink.png",
        "ANDREA (1:00): Concretely, Brink is four things. A feed of songs that play right there in the card, "
        "with reactions and comments. A taste community for each person — a group our model puts you in. A "
        "compatibility score between any two people. And a small artist portal. It's all one live app, and "
        "everything we set out to build is built. But the part we actually want to talk about is how three "
        "of us built it.",
    )

    # 4 — What we set out to do (the AI narrative)
    add_content(
        prs, "ANDREA · 1:30", "What we really set out to do",
        ["A big goal was to see how far we could get by leaning on AI coding agents.",
         "We learned fast: **AI is only as good as the context you give it.",
         "So the job was two things — give the agents rich context, and be exact about what to build.",
         "**Our first couple of weeks were planning, not code:  ADRs, tickets, and project governance.",
         "That groundwork is what let everything after it move quickly."],
        "ANDREA (1:30): Honestly, one of our real goals on this project was to see how far we could get by "
        "leaning on AI coding agents — to build as much as we possibly could as a team of three. And what "
        "we learned early is that AI is only as good as the context you give it. If you're vague, you get "
        "vague. So the whole job became two things: give the agents as much context as possible, and be "
        "really precise about exactly what we wanted built. That's why our first couple of weeks weren't "
        "spent writing app code at all — we were setting up what we call our ADRs, our tickets, and our "
        "project governance. That planning is what made everything afterwards move fast.",
    )

    # 5 — The scaffolding (three cards)
    add_columns(
        prs, "ANDREA · 1:30", "The scaffolding that let AI build",
        [("ADRs", ["Every real decision, written down with the reasoning.",
                   "So an agent — or a teammate — knows why things are the way they are.",
                   "Append-only: we supersede, never quietly edit."], ACCENT),
         ("Tickets", ["One clear piece of work at a time.",
                      "The agent builds exactly that — nothing more.",
                      "Each one traces back to a requirement."], HOT),
         ("Governance", ["Guardrails around every change: one PR at a time.",
                         "CI runs tests + a secret scan automatically.",
                         "A check fails if you change code but not its docs."], TEAL)],
        "ANDREA (1:30): So what does that scaffolding actually look like. ADRs are architecture decision "
        "records — every real decision written down with the reasoning behind it, so anyone picking up the "
        "work, an agent or one of us, knows why things are the way they are. Tickets break the work into "
        "one clear piece at a time, so the agent builds exactly that and nothing else. And governance is "
        "the guardrails — every change is a pull request, tests and a secret-scanner run automatically, and "
        "we even have a check that fails if you change code without updating its docs. All of that is "
        "context and clarity — which is exactly what the agents needed to be useful.",
    )

    # 6 — Two pivots
    add_content(
        prs, "ANDREA · 1:00", "Two times we changed our minds",
        ["Backend:  TypeScript / Vercel  →  Python / FastAPI   (ADR-0010)",
         "Frontend:  React app  →  simpler server-rendered pages   (ADR-0013)",
         "**Why both times: the security-sensitive parts were in a stack we couldn't all read.",
         "Now it's one Python codebase all three of us can actually read and defend."],
        "ANDREA (1:00): There were two moments where we changed our minds, and they're worth sharing "
        "because the reasoning was the same both times. We started the backend in TypeScript and the "
        "frontend as a separate React app. But the security-sensitive parts — logins, encrypting people's "
        "tokens — were in a stack not all three of us could confidently read. For a graded project, that's "
        "a real risk. So we moved the backend to Python, and later replaced the React app with simpler "
        "server-rendered pages. Now it's one Python codebase all of us can read and defend. Both of those "
        "are written up as ADRs, so you can see the reasoning.",
    )

    # 7 — Architecture diagram
    add_architecture(
        prs, "ANDREA · 1:00", "How it fits together",
        "ANDREA (1:00): Here's how it all fits together. Everything the user touches is one app — a FastAPI "
        "service on Render that serves both the web pages and the API, so there's no separate frontend to "
        "deploy. It talks to Supabase for the database, logins, and file storage. Off to the side, a Python "
        "job runs every night on GitHub Actions, trains our model, and writes the results into the "
        "database — and the app just reads those results when someone loads a page. And Spotify is where the "
        "logins and the listening data come from.",
    )

    # 8 — Data model
    add_content(
        prs, "ANDREA · 0:50", "The data behind it  (~14 tables)",
        ["**Social:  users, posts, reactions, comments, the follow graph",
         "**Music & model:  tracks, plays, and the taste communities the model produces",
         "A person's community and compatibility aren't stored — we work them out on the spot,",
         "so they're never stale and always reflect the latest listening."],
        "ANDREA (0:50): Behind it is about 14 tables. The social side is what you'd expect — users, posts, "
        "reactions, comments, who follows whom. The music side holds the tracks, everyone's plays, and the "
        "taste communities the model produces. One thing we're a little proud of: a person's community and "
        "their compatibility with you aren't saved anywhere — we work them out on the spot when you load "
        "the page, so they're never out of date.",
    )

    # 9 — Auth & security
    add_content(
        prs, "ANDREA · 0:45", "Keeping accounts safe",
        ["We check every login on our server — we never hold a master key.",
         "**Spotify tokens are encrypted before they're stored  (AES-256-GCM).",
         "Limits on how fast anyone can post, and a secret-scanner on every commit.",
         "Security was the reason we moved to Python — so we took it seriously."],
        "ANDREA (0:45): On security — we validate every login on our own server and never hold a master "
        "secret. The tokens we keep for Spotify are encrypted before they ever touch the database. There "
        "are limits on how fast anyone can post, and a scanner runs on every single commit to make sure we "
        "never accidentally check in a password or key. Security was the reason we moved to Python in the "
        "first place, so we didn't cut corners here. I'll hand over to Jonah for the analytics.",
    )

    # 10 — Analytics: the data layers
    add_content(
        prs, "JONAH · 1:15", "How the data is organised",
        ["**Raw — everything exactly as it arrives (Spotify pulls, dataset rows)",
         "**Cleaned — tidied into the tables the app actually uses",
         "**Model output — the communities and the trained model itself",
         "It's a standard 'medallion' pattern. A bit more structure than we strictly needed —",
         "but everything stays traceable, which is the point."],
        "JONAH (1:15): The data side is organised in three layers — people call it a medallion. The first "
        "layer is raw: everything exactly as it arrives, the Spotify listening pulls and the raw dataset "
        "rows. The second is that data cleaned up into the tables the app actually uses. And the third is "
        "the output of our model — the communities and the model itself. It's a common data-engineering "
        "pattern, and honestly at our scale it's a little more structure than we strictly needed — but it "
        "keeps everything traceable, and it's the right way to do it.",
    )

    # 11 — Analytics: the model (split + analytics screenshot)
    add_split(
        prs, "JONAH · 1:30", "The model, in two parts",
        ["**Once a night:  we cluster ~1.2M songs by their sound",
         "(danceable, energetic, upbeat…) into 7 taste communities.",
         "**When you open a profile:  we place that person's songs against",
         "the communities and find where they fit.",
         "Compatibility = how close two people's taste profiles are.",
         "Nothing's hard-coded — retrain tonight, the app updates tomorrow."],
        "analytics-brink.png",
        "JONAH (1:30): The model is in two parts. Once a night, we train a clustering model on about 1.2 "
        "million songs, using audio features like how danceable or energetic or upbeat a track is — that "
        "sorts songs into seven taste communities. Then, when you actually open someone's profile, we take "
        "the songs they've played, place them against those communities, and find where they fit. "
        "Compatibility between two people is just how close their two taste profiles are. And nothing is "
        "hard-coded — if we retrain the model tonight, the app picks up the new version tomorrow. This page "
        "is that model's output, live.",
    )

    # 12 — Analytics: honesty
    add_content(
        prs, "JONAH · 0:45", "Being honest about the limits",
        ["The math actually liked just 2 groups — we chose 7 on purpose, and we say so.",
         "Only ~a quarter of songs matched our dataset — we report that, we don't hide it.",
         "**We planned a second model (song popularity) and cut it — the data couldn't support it.",
         "We wrote down why. We'd rather show the real picture than a polished one."],
        "JONAH (0:45): We also want to be upfront about the limits, because we think that matters. The math "
        "actually suggested only two groups was cleanest — but two isn't very interesting as 'communities,' "
        "so we chose seven on purpose, and we say that openly. Only about a quarter of songs matched our "
        "dataset, and we report that honestly rather than hiding it. And we'd planned a second model to "
        "predict song popularity, but the data just couldn't support it honestly, so we cut it — and wrote "
        "down why. We'd rather show you the real picture than a polished one. Sebastian will show where all "
        "this shows up in the app.",
    )

    # 13 — Frontend (split + landing screenshot)
    add_split(
        prs, "SEBASTIAN · 1:00", "One look, five pages",
        ["Dark, calm, music-first — a lavender & pink accent throughout.",
         "**Landing · Feed · Profile · Artist studio · Analytics",
         "No separate app to build — the pages come from the same server,",
         "with light JavaScript for the interactive bits.",
         "That kept it fast, and simple for a small team to maintain."],
        "Landing-brink.png",
        "SEBASTIAN (1:00): The whole frontend has one consistent look — dark, calm, music-first, with a "
        "lavender-and-pink accent. There are five pages: the landing page, the feed, profiles, an artist "
        "studio, and the analytics page. There's no separate app to build — the pages come straight from "
        "the same server, with just a bit of JavaScript for the interactive parts. That kept everything "
        "fast to load and simple for a small team to keep up with.",
    )

    # 14 — Frontend details (split + profile screenshot)
    add_split(
        prs, "SEBASTIAN · 0:45", "The small touches",
        ["**Tap the album art — a Spotify player opens right in the card.",
         "**A 'liked by' line, and double-tap to heart, like you'd expect.",
         "**One tap to share whatever you're playing right now.",
         "Every empty screen has a friendly nudge, not a blank space."],
        "Profile-brink.png",
        "SEBASTIAN (0:45): The small touches are what make it feel like a real social app. Tap the album "
        "art and a Spotify player opens right there in the card. There's a 'liked by' line, and you can "
        "double-tap to heart a song like you'd expect. There's a one-tap button to share whatever you're "
        "playing right now. And every empty screen — no posts yet, nothing playing — has a friendly nudge "
        "instead of a blank space. Let me just show you the real thing.",
    )

    # 15 — Live demo
    add_content(
        prs, "SEBASTIAN drives · Jonah + Andrea narrate · 2:30", "Let's look at the real thing",
        ["brink-xg7p.onrender.com",
         "1.  Feed — scroll, tap the art to play, react and comment",
         "2.  Share — post the track I'm listening to right now",
         "3.  Profile — taste community + how compatible we are   (Andrea)",
         "4.  Analytics — the 7 communities and their 'audio DNA'   (Jonah)",
         "**Warm the site up first. The next slides are screenshots, as a backup."],
        "SEBASTIAN drives (~2:30). Pre-warm the site so there's no wait. 1) Feed: 'this is the live app' — "
        "scroll, tap album art so a song plays in the card, react and drop a comment. 2) Share: use the "
        "one-tap 'share what you're playing' and post it. 3) Profile (Andrea narrates): open someone "
        "else's profile — their streak, top tracks, their taste community, and since it's not me, a "
        "compatibility score, all worked out live. 4) Analytics (Jonah narrates): the seven communities, "
        "how distinct they are, and each one's audio DNA — real numbers from last night's run. "
        "IF ANYTHING BREAKS: switch to the screenshots on the next slides.",
    )

    # Fallback screenshots
    for img, cap in [
        ("Landing-brink.png", "Landing — sign in; your listening becomes the feed"),
        ("Feed-brink.png", "Feed — songs play in place, with reactions & comments"),
        ("Profile-brink.png", "Profile — taste community + how compatible you are"),
        ("artist-brink.png", "Artist studio — behind-the-scenes posts"),
        ("analytics-brink.png", "Analytics — the model's real output: communities & audio DNA"),
    ]:
        p = SHOTS / img
        if p.exists():
            add_image(prs, p, cap, "Demo backup / reference: " + cap)

    # 16 — What we learned
    add_content(
        prs, "ALL · 0:45", "What we took away",
        ["**The AI was only as good as the context and structure we gave it.  (Andrea)",
         "**Being honest about the model's limits made it more convincing, not less.  (Jonah)",
         "**One codebase and one look let three people move without stepping on each other.  (Sebastian)"],
        "ALL (0:45), one line each. ANDREA: the biggest lesson for me was that the AI was only ever as good "
        "as the context and structure we gave it — the planning up front was the whole game. JONAH: being "
        "honest about the model's limits made it more convincing, not less. SEBASTIAN: keeping it to one "
        "codebase and one consistent look is what let three of us move fast without stepping on each other.",
    )

    # 17 — Close
    add_title(
        prs, "Brink — Music, but social. 🎶", "Thanks — happy to take questions.",
        ["brink-xg7p.onrender.com"],
        "ANDREA (0:15): That's Brink — live, and built by the three of us. Thanks for listening — we're "
        "happy to take any questions.",
    )

    import os
    out = Path(os.environ.get("BRINK_DECK_OUT", Path(__file__).with_name("brink-deck.pptx")))
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"Wrote {out}  ({n} slides)")


if __name__ == "__main__":
    build()
